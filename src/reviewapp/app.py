"""이커머스 리뷰/데이터 분석 웹앱 - FastAPI 메인."""
import os
import re
import uuid
import shutil
import json
import logging
import traceback

from dotenv import load_dotenv
load_dotenv()

import httpx
from fastapi import FastAPI, UploadFile, File, Form, Request, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse, FileResponse
from fastapi.staticfiles import StaticFiles
from fastapi.templating import Jinja2Templates

import pandas as pd

# ---------------------------------------------------------------------------
# 앱 설정
# ---------------------------------------------------------------------------
PKG_DIR = os.path.dirname(os.path.abspath(__file__))   # 패키지 리소스 (templates, static, fonts, data)
WORK_DIR = os.path.join(os.path.expanduser("~"), ".review-analyzer")  # 사용자 데이터
UPLOAD_DIR = os.path.join(WORK_DIR, "uploads")
RESULT_DIR = os.path.join(WORK_DIR, "results")
os.makedirs(UPLOAD_DIR, exist_ok=True)
os.makedirs(RESULT_DIR, exist_ok=True)

logging.basicConfig(level=logging.INFO, format="%(asctime)s %(levelname)s %(message)s")
logger = logging.getLogger("review-analyzer")

INSIGHT_MODEL = "anthropic/claude-sonnet-4-6"

app = FastAPI(title="이커머스 데이터 분석", description="수업용 리뷰/데이터 분석 웹앱")
app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_methods=["*"], allow_headers=["*"])
app.mount("/static", StaticFiles(directory=os.path.join(PKG_DIR, "static")), name="static")
templates = Jinja2Templates(directory=os.path.join(PKG_DIR, "templates"))

# ---------------------------------------------------------------------------
# 공용 유틸
# ---------------------------------------------------------------------------

def _extract_text_lines(save_path: str, ext: str) -> list[str]:
    """비정형 파일에서 텍스트 라인 추출."""
    lines = []
    if ext == ".txt":
        with open(save_path, "r", encoding="utf-8", errors="ignore") as f:
            lines = [line.strip() for line in f if line.strip()]
    elif ext == ".pdf":
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(save_path)
            for page in doc:
                text = page.get_text()
                lines.extend([l.strip() for l in text.split("\n") if l.strip()])
            doc.close()
        except ImportError:
            raise RuntimeError("PDF 파싱에 PyMuPDF가 필요합니다. pip install pymupdf")
    elif ext in (".doc", ".docx"):
        try:
            import docx
            doc = docx.Document(save_path)
            lines = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        except ImportError:
            raise RuntimeError("DOC/DOCX 파싱에 python-docx가 필요합니다. pip install python-docx")
    elif ext in (".ppt", ".pptx"):
        try:
            from pptx import Presentation
            prs = Presentation(save_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                lines.append(text)
        except ImportError:
            raise RuntimeError("PPT/PPTX 파싱에 python-pptx가 필요합니다. pip install python-pptx")
    return lines


def load_dataframe(file: UploadFile) -> tuple[pd.DataFrame, str]:
    """업로드 파일 → DataFrame + 저장 경로."""
    job_id = str(uuid.uuid4())[:8]
    job_dir = os.path.join(RESULT_DIR, job_id)
    os.makedirs(job_dir, exist_ok=True)
    ext = os.path.splitext(file.filename)[1].lower()
    save_path = os.path.join(UPLOAD_DIR, f"{job_id}{ext}")
    with open(save_path, "wb") as f:
        shutil.copyfileobj(file.file, f)

    if ext == ".csv":
        df = pd.read_csv(save_path)
    elif ext in (".xlsx", ".xls"):
        df = pd.read_excel(save_path)
    elif ext in (".txt", ".pdf", ".doc", ".docx", ".ppt", ".pptx"):
        lines = _extract_text_lines(save_path, ext)
        if not lines:
            raise ValueError("파일에서 텍스트를 추출할 수 없습니다.")
        df = pd.DataFrame({"텍스트": lines})
    else:
        raise ValueError(f"지원하지 않는 파일 형식입니다: {ext}")

    return df, job_id


def detect_text_column(df: pd.DataFrame) -> str:
    """리뷰/텍스트 컬럼 자동 탐지."""
    priority = ["리뷰", "내용", "텍스트", "후기", "comment", "review", "content", "body", "text"]
    candidates = []
    for col in df.columns:
        sample = df[col].dropna().astype(str)
        if len(sample) == 0:
            continue
        avg_len = sample.str.len().mean()
        candidates.append((col, avg_len))
    candidates.sort(key=lambda x: x[1], reverse=True)
    for col, _ in candidates:
        if any(kw in col.lower() for kw in priority):
            return col
    return candidates[0][0] if candidates else df.columns[0]


# ---------------------------------------------------------------------------
# 페이지 라우트
# ---------------------------------------------------------------------------

@app.get("/")
async def index(request: Request):
    return templates.TemplateResponse("index.html", {"request": request})


# ---------------------------------------------------------------------------
# API 엔드포인트 - 탭1: 리뷰 분류
# ---------------------------------------------------------------------------

@app.post("/api/classify")
async def api_classify(file: UploadFile = File(...), chart_mode: str = Form("plotly")):
    try:
        from reviewapp.analyzer.classify import run_classification
        df, job_id = load_dataframe(file)
        result = run_classification(df, job_id, chart_mode)
        _save_insight_context(job_id, "classify", result)
        return JSONResponse({"status": "ok", "job_id": job_id, **result})
    except Exception as e:
        logger.error(traceback.format_exc())
        return JSONResponse({"status": "error", "message": str(e)}, status_code=500)


# ---------------------------------------------------------------------------
# API 엔드포인트 - 탭2: 감성 분석
# ---------------------------------------------------------------------------

@app.post("/api/sentiment")
async def api_sentiment(file: UploadFile = File(...), chart_mode: str = Form("plotly")):
    try:
        from reviewapp.analyzer.sentiment import run_sentiment
        df, job_id = load_dataframe(file)
        result = run_sentiment(df, job_id, chart_mode)
        _save_insight_context(job_id, "sentiment", result)
        return JSONResponse({"status": "ok", "job_id": job_id, **result})
    except Exception as e:
        logger.error(traceback.format_exc())
        return JSONResponse({"status": "error", "message": str(e)}, status_code=500)


# ---------------------------------------------------------------------------
# API 엔드포인트 - 탭3: EDA
# ---------------------------------------------------------------------------

@app.post("/api/eda")
async def api_eda(file: UploadFile = File(...), chart_mode: str = Form("plotly")):
    try:
        from reviewapp.analyzer.eda import run_eda
        df, job_id = load_dataframe(file)
        result = run_eda(df, job_id, chart_mode)
        _save_insight_context(job_id, "eda", result)
        return JSONResponse({"status": "ok", "job_id": job_id, **result})
    except Exception as e:
        logger.error(traceback.format_exc())
        return JSONResponse({"status": "error", "message": str(e)}, status_code=500)


# ---------------------------------------------------------------------------
# API 엔드포인트 - 탭4: 커머스 분석
# ---------------------------------------------------------------------------

@app.post("/api/tcp")
async def api_tcp(
    file: UploadFile = File(...),
    chart_mode: str = Form("plotly"),
    dimensions: str = Form("time,customer,product"),
    date_col: str = Form(""),
    customer_col: str = Form(""),
    product_col: str = Form(""),
    quantity_col: str = Form(""),
    price_col: str = Form(""),
):
    try:
        from reviewapp.analyzer.tcp_rfm import run_tcp
        df, job_id = load_dataframe(file)
        col_map = {
            "date": date_col, "customer": customer_col,
            "product": product_col, "quantity": quantity_col, "price": price_col,
        }
        result = run_tcp(df, job_id, chart_mode, col_map, dimensions=dimensions.split(","))
        _save_insight_context(job_id, "tcp", result)
        return JSONResponse({"status": "ok", "job_id": job_id, **result})
    except Exception as e:
        logger.error(traceback.format_exc())
        return JSONResponse({"status": "error", "message": str(e)}, status_code=500)


# ---------------------------------------------------------------------------
# API 엔드포인트 - 탭5: 텍스트 마이닝
# ---------------------------------------------------------------------------

@app.post("/api/textmining")
async def api_textmining(
    file: UploadFile = File(...),
    chart_mode: str = Form("plotly"),
    n_topics: int = Form(5),
    top_n: int = Form(20),
    pos_filter: str = Form("NNG,NNP"),
    analyses: str = Form("tfidf,lda,network,wordcloud,sentiment"),
):
    try:
        from reviewapp.analyzer.text_mining import run_textmining
        df, job_id = load_dataframe(file)
        result = run_textmining(
            df, job_id, chart_mode,
            n_topics=n_topics, top_n=top_n,
            pos_filter=pos_filter.split(","),
            analyses=analyses.split(","),
        )
        _save_insight_context(job_id, "textmining", result)
        return JSONResponse({"status": "ok", "job_id": job_id, **result})
    except Exception as e:
        logger.error(traceback.format_exc())
        return JSONResponse({"status": "error", "message": str(e)}, status_code=500)


# ---------------------------------------------------------------------------
# 다운로드 엔드포인트
# ---------------------------------------------------------------------------

@app.get("/api/download/{job_id}/{filename}")
async def download_file(job_id: str, filename: str):
    path = os.path.join(RESULT_DIR, job_id, filename)
    if not os.path.exists(path):
        raise HTTPException(404, "파일을 찾을 수 없습니다")
    return FileResponse(path, filename=filename)


# ---------------------------------------------------------------------------
# AI 인사이트 - 컨텍스트 저장 및 채팅
# ---------------------------------------------------------------------------

ANALYSIS_NAMES = {
    "classify": "리뷰 분류", "sentiment": "감성 분석", "eda": "EDA 분석",
    "tcp": "커머스 분석", "textmining": "텍스트 마이닝",
}

def _strip_html(html: str) -> str:
    """HTML 태그 제거하여 텍스트만 추출."""
    text = re.sub(r'<[^>]+>', ' ', html)
    return re.sub(r'\s+', ' ', text).strip()


def _save_insight_context(job_id: str, analysis_type: str, result: dict):
    """분석 결과를 텍스트로 저장하여 AI 인사이트 컨텍스트로 사용."""
    parts = [f"분석 유형: {ANALYSIS_NAMES.get(analysis_type, analysis_type)}"]
    if result.get("summary_html"):
        parts.append(f"분석 요약:\n{_strip_html(result['summary_html'])}")
    if result.get("details_html"):
        parts.append(f"상세 결과:\n{_strip_html(result['details_html'])}")
    context = "\n\n".join(parts)
    path = os.path.join(RESULT_DIR, job_id, "insight_context.txt")
    with open(path, "w", encoding="utf-8") as f:
        f.write(context)


@app.post("/api/insight/init")
async def api_insight_init(request: Request):
    """분석 결과 초기 AI 해석 생성."""
    body = await request.json()
    job_id = body.get("job_id")
    api_key = body.get("api_key", "")
    if not job_id or not api_key:
        return JSONResponse({"status": "error", "message": "API 키 또는 job_id 없음"}, status_code=400)

    context_path = os.path.join(RESULT_DIR, job_id, "insight_context.txt")
    if not os.path.exists(context_path):
        return JSONResponse({"status": "error", "message": "분석 결과를 찾을 수 없습니다"}, status_code=404)

    with open(context_path, "r", encoding="utf-8") as f:
        context = f.read()

    messages = [
        {"role": "system", "content": (
            "당신은 데이터 분석 전문가입니다. 아래 분석 결과를 바탕으로 핵심 인사이트를 한국어로 요약해주세요.\n"
            "- 3~5개 핵심 포인트를 간결하게 정리\n"
            "- 수치를 근거로 들어 설명\n"
            "- 데이터에 없는 내용은 추측하지 마세요\n"
            "- 이모지를 사용하지 마세요\n"
            "- 마크다운은 ##, **, - 만 사용하세요\n\n"
            f"분석 결과:\n{context}"
        )},
        {"role": "user", "content": "이 분석 결과의 핵심 인사이트를 알려주세요."},
    ]

    try:
        async with httpx.AsyncClient() as client:
            resp = await client.post(
                "https://openrouter.ai/api/v1/chat/completions",
                headers={"Authorization": f"Bearer {api_key}"},
                json={"model": INSIGHT_MODEL, "messages": messages, "max_tokens": 10284},
                timeout=60,
            )
            data = resp.json()
            if "error" in data:
                return JSONResponse({"status": "error", "message": data["error"].get("message", "API 오류")}, status_code=401)
            answer = data["choices"][0]["message"]["content"]
        return JSONResponse({"status": "ok", "message": answer})
    except Exception as e:
        logger.error(f"Insight init error: {e}")
        return JSONResponse({"status": "error", "message": str(e)}, status_code=500)


@app.post("/api/insight/chat")
async def api_insight_chat(request: Request):
    """추가 질문 응답."""
    body = await request.json()
    job_id = body.get("job_id")
    api_key = body.get("api_key", "")
    question = body.get("message", "")
    history = body.get("history", [])

    if not job_id or not question or not api_key:
        return JSONResponse({"status": "error", "message": "필수 파라미터 누락"}, status_code=400)

    context_path = os.path.join(RESULT_DIR, job_id, "insight_context.txt")
    if not os.path.exists(context_path):
        return JSONResponse({"status": "error", "message": "분석 결과를 찾을 수 없습니다"}, status_code=404)

    with open(context_path, "r", encoding="utf-8") as f:
        context = f.read()

    messages = [
        {"role": "system", "content": (
            "당신은 데이터 분석 전문가입니다. 아래 분석 결과를 바탕으로 사용자의 질문에 한국어로 답변해주세요.\n"
            "- 수치를 근거로 들어 설명\n"
            "- 데이터에 없는 내용은 추측하지 마세요\n"
            "- 간결하고 명확하게 답변하세요\n"
            "- 이모지를 사용하지 마세요\n"
            "- 마크다운은 ##, **, - 만 사용하세요\n\n"
            f"분석 결과:\n{context}"
        )},
    ]
    for h in history[-10:]:
        messages.append({"role": h["role"], "content": h["content"]})
    messages.append({"role": "user", "content": question})

    try:
        async with httpx.AsyncClient() as client:
            resp = await client.post(
                "https://openrouter.ai/api/v1/chat/completions",
                headers={"Authorization": f"Bearer {api_key}"},
                json={"model": INSIGHT_MODEL, "messages": messages, "max_tokens": 10284},
                timeout=60,
            )
            data = resp.json()
            if "error" in data:
                return JSONResponse({"status": "error", "message": data["error"].get("message", "API 오류")}, status_code=401)
            answer = data["choices"][0]["message"]["content"]
        return JSONResponse({"status": "ok", "message": answer})
    except Exception as e:
        logger.error(f"Insight chat error: {e}")
        return JSONResponse({"status": "error", "message": str(e)}, status_code=500)


# ---------------------------------------------------------------------------
# 헬스체크
# ---------------------------------------------------------------------------

@app.get("/health")
async def health():
    return {"status": "ok"}


if __name__ == "__main__":
    import uvicorn
    uvicorn.run("app:app", host="0.0.0.0", port=int(os.getenv("PORT", 8000)), reload=True)
